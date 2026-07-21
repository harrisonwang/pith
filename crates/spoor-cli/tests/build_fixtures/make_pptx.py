#!/usr/bin/env python3
"""PPTX fixtures."""
import base64
from io import BytesIO
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

OUT = Path(__file__).resolve().parent.parent / "fixtures" / "pptx"
OUT.mkdir(parents=True, exist_ok=True)


# ---------- 01: basic title slides ----------
def build_01_basic():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Title Slide"
    s.placeholders[1].text = "A subtitle"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Second slide"
    tf = s.placeholders[1].text_frame
    tf.text = "First bullet"
    p = tf.add_paragraph(); p.text = "Second bullet"
    p = tf.add_paragraph(); p.text = "Nested item"; p.level = 1

    prs.save(OUT / "01_basic.pptx")


# ---------- 02: with table ----------
def build_02_with_table():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Has a Table"
    rows, cols = 3, 2
    tbl = s.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(4), Inches(2)).table
    tbl.cell(0, 0).text = "Header A"
    tbl.cell(0, 1).text = "Header B"
    tbl.cell(1, 0).text = "1"
    tbl.cell(1, 1).text = "2"
    tbl.cell(2, 0).text = "3"
    tbl.cell(2, 1).text = "4"
    prs.save(OUT / "02_with_table.pptx")


# ---------- 03: speaker notes ----------
def build_03_with_notes():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide with notes"
    s.notes_slide.notes_text_frame.text = "These are speaker notes that explain the slide."
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide without notes"
    prs.save(OUT / "03_with_notes.pptx")


# ---------- 04: empty deck (one blank slide) ----------
def build_04_empty():
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    prs.save(OUT / "04_empty.pptx")


# ---------- 05: many slides (ordering test - ensure slide11 > slide2) ----------
def build_05_ordering():
    prs = Presentation()
    for i in range(1, 13):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"Slide number {i}"
    prs.save(OUT / "05_ordering.pptx")


# ---------- 06: merged table must emit an integrity warning ----------
def build_06_merged_table():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Merged table"
    tbl = s.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(2)).table
    tbl.cell(0, 0).text = "Merged header"
    tbl.cell(0, 0).merge(tbl.cell(0, 1))
    tbl.cell(1, 0).text = "A"
    tbl.cell(1, 1).text = "B"
    prs.save(OUT / "06_merged_table.pptx")


# ---------- 07: omitted picture must emit an integrity warning ----------
def build_07_embedded_visual():
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
        "+A8AAQUBAScY42YAAAAASUVORK5CYII="
    )
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Picture slide"
    s.shapes.add_picture(BytesIO(png), Inches(1), Inches(2), Inches(2), Inches(2))
    prs.save(OUT / "07_embedded_visual.pptx")


# ---------- 08: multi-slide image placeholders ----------
# Drives the `spoor://pptx/part/ppt/media/*` emission path: covers per-slide
# image numbering, multiple images on a single slide, and an image-free slide
# that should produce no handles.
def build_08_image_placeholders():
    # Two distinct 1x1 PNGs so python-pptx writes two separate media parts
    # instead of deduplicating by content hash.
    red_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8/5+hHgAHggJ/PchI7wAAAABJRU5ErkJggg=="
    )
    blue_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNgYPj/HwADAgH/eL5gtgAAAABJRU5ErkJggg=="
    )
    prs = Presentation()
    # Slide 1: one image.
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide one"
    s.shapes.add_picture(BytesIO(red_png), Inches(1), Inches(2), Inches(2), Inches(2))
    # Slide 2: two images, distinct bytes so they land in distinct media parts.
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide two"
    s.shapes.add_picture(BytesIO(red_png), Inches(1), Inches(2), Inches(2), Inches(2))
    s.shapes.add_picture(BytesIO(blue_png), Inches(4), Inches(2), Inches(2), Inches(2))
    # Slide 3: title only — must produce no handle.
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide three (no images)"
    prs.save(OUT / "08_image_placeholders.pptx")


# ---------- 09: sldIdLst order differs from filename numbers ----------
# python-pptx appends parts in creation order (slide1=Alpha, slide2=Beta,
# slide3=Gamma); we then move Gamma to the front of sldIdLst so the deck
# plays Gamma, Alpha, Beta. Slide numbers must follow deck order, not the
# part filenames.
def build_09_reordered():
    prs = Presentation()
    for title in ("Alpha", "Beta", "Gamma"):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = title
    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst)
    sld_id_lst.insert(0, ids[2])  # lxml moves the element
    prs.save(OUT / "09_reordered.pptx")


# ---------- 10: hidden slide keeps its number ----------
# Slide 2 carries show="0" (PowerPoint's "Hide Slide"). Its body must be
# omitted while its number is kept, so numbering stays aligned with what
# PowerPoint displays.
def build_10_hidden_slide():
    prs = Presentation()
    for title in ("Visible one", "Secret draft", "Visible three"):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = title
    prs.slides[1]._element.set("show", "0")
    prs.save(OUT / "10_hidden_slide.pptx")


# ---------- 11: image-only slides (no text layer) ----------
# Slide 1: a picture and speaker notes, no body text. Slide 2: picture only.
# Slide 3: a text slide as control. Drives the slide_no_text_layer posture
# warning and its notes-aware wording.
def build_11_image_only():
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
        "+A8AAQUBAScY42YAAAAASUVORK5CYII="
    )
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(BytesIO(png), Inches(1), Inches(1), Inches(8), Inches(5))
    s.notes_slide.notes_text_frame.text = "The whole slide is a screenshot; the script lives here."
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(BytesIO(png), Inches(2), Inches(2), Inches(4), Inches(3))
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Text control slide"
    # Slide 4: the real-world "pure image" shape — a title placeholder plus a
    # full-bleed screenshot. The title is a label, not a text layer.
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Architecture overview"
    s.shapes.add_picture(BytesIO(png), Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    prs.save(OUT / "11_image_only.pptx")


# ---------- 12: reading order beats XML (z-) order ----------
# Text boxes added bottom-first: XML order is Bottom row, Top right, Top left;
# the geometric (top, left) sort must output Top left, Top right, Bottom row.
def build_12_reading_order():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    b = s.shapes.add_textbox(Inches(1), Inches(5), Inches(3), Inches(1))
    b.text_frame.text = "Bottom row"
    r = s.shapes.add_textbox(Inches(6), Inches(2), Inches(3), Inches(1))
    r.text_frame.text = "Top right"
    t = s.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    t.text_frame.text = "Top left"
    prs.save(OUT / "12_reading_order.pptx")


# ---------- 13: bullet levels, numbering and opt-out ----------
def build_13_bullets():
    from pptx.oxml.ns import qn

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Bulleted body"
    tf = s.placeholders[1].text_frame
    tf.text = "Level zero"
    p = tf.add_paragraph(); p.text = "Level one"; p.level = 1
    p = tf.add_paragraph(); p.text = "Level two"; p.level = 2
    p = tf.add_paragraph(); p.text = "Numbered item"
    pPr = p._p.get_or_add_pPr()
    pPr.append(pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"}))
    p = tf.add_paragraph(); p.text = "No bullet here"
    pPr = p._p.get_or_add_pPr()
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))
    prs.save(OUT / "13_bullets.pptx")


# ---------- 14: author alt text on a picture, with injection attempt ----------
def build_14_alt_text():
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
        "+A8AAQUBAScY42YAAAAASUVORK5CYII="
    )
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    pic = s.shapes.add_picture(BytesIO(png), Inches(1), Inches(1), Inches(2), Inches(2))
    pic._element._nvXxPr.cNvPr.set("descr", "Quarterly revenue chart ](evil) [x]")
    prs.save(OUT / "14_alt_text.pptx")


# ---------- 15: group shapes flatten in order ----------
def build_15_group_shapes():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    g = s.shapes.add_group_shape()
    a = g.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    a.text_frame.text = "Grouped alpha"
    b = g.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    b.text_frame.text = "Grouped beta"
    lone = s.shapes.add_textbox(Inches(1), Inches(4), Inches(2), Inches(1))
    lone.text_frame.text = "After the group"
    prs.save(OUT / "15_group_shapes.pptx")


# ---------- 16: notes furniture (page number) must not leak ----------
def build_16_notes_furniture():
    import copy
    from pptx.oxml.ns import qn

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Notes with furniture"
    ns = s.notes_slide
    ns.notes_text_frame.text = "Real speaker notes."
    body_sp = ns.notes_text_frame._txBody.getparent()
    dup = copy.deepcopy(body_sp)
    ph = dup.find(".//" + qn("p:ph"))
    ph.set("type", "sldNum")
    ph.set("idx", "10")
    for t in dup.findall(".//" + qn("a:t")):
        t.text = "12"
    body_sp.getparent().append(dup)
    prs.save(OUT / "16_notes_furniture.pptx")


# ---------- 17: no presentation.xml → filename-order fallback ----------
# A hand-degraded package (presentation part stripped) must fall back to
# numeric filename order deterministically instead of erroring or guessing.
def build_17_no_presentation():
    import zipfile

    src = OUT / "09_reordered.pptx"
    dst = OUT / "17_no_presentation.pptx"
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w") as zout:
        for item in zin.namelist():
            if item == "ppt/presentation.xml":
                continue
            zout.writestr(item, zin.read(item))


# ---------- 18: native chart data ----------
# Business decks put their key numbers in charts; the cached c:cat/c:val data
# must come out as a table, and a fully-extracted chart slide must not claim
# its output is incomplete.
def build_18_chart():
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Numbers live in the chart"
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("Revenue", (100, 120, 140))
    data.add_series("Cost", (80, 90, 95))
    frame = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(8), Inches(4), data
    )
    frame.chart.has_title = True
    frame.chart.chart_title.text_frame.text = "Quarterly performance"
    prs.save(OUT / "18_chart.pptx")


# ---------- 19: SmartArt node text (zip surgery; python-pptx can't author it) ----------
def build_19_smartart():
    import zipfile

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    tmp = OUT / "_19_base.pptx"
    prs.save(tmp)

    NS = (
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
    )
    frame = (
        '<p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id="7" name="SmartArt"/>'
        "<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>"
        '<p:xfrm><a:off x="914400" y="914400"/><a:ext cx="6096000" cy="4572000"/></p:xfrm>'
        '<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
        '<dgm:relIds r:dm="rId80" r:lo="rId81" r:qs="rId82" r:cs="rId83"/>'
        "</a:graphicData></a:graphic></p:graphicFrame>"
    )

    def pt(model_id, text=None, pt_type=None):
        t = f' type="{pt_type}"' if pt_type else ""
        body = (
            f"<dgm:t><a:bodyPr/><a:p><a:r><a:t>{text}</a:t></a:r></a:p></dgm:t>" if text else ""
        )
        return f'<dgm:pt modelId="{model_id}"{t}><dgm:prSet/><dgm:spPr/>{body}</dgm:pt>'

    data_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f"<dgm:dataModel {NS}><dgm:ptLst>"
        + pt("{D0}", pt_type="doc")
        + pt("{N1}", "Plan")
        + pt("{T1}", pt_type="parTrans")
        + pt("{N2}", "Build")
        + pt("{N3}", "Ship")
        + "</dgm:ptLst><dgm:bg/><dgm:whole/></dgm:dataModel>"
    )

    rel = (
        '<Relationship Id="rId80" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData" '
        'Target="../diagrams/data1.xml"/>'
    )
    override = (
        '<Override PartName="/ppt/diagrams/data1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.drawingml.diagramData+xml"/>'
    )

    with zipfile.ZipFile(tmp) as zin, zipfile.ZipFile(OUT / "19_smartart.pptx", "w") as zout:
        for item in zin.namelist():
            content = zin.read(item)
            if item == "ppt/slides/slide1.xml":
                text = content.decode()
                text = text.replace("</p:spTree>", frame + "</p:spTree>")
                content = text.encode()
            elif item == "ppt/slides/_rels/slide1.xml.rels":
                content = zin.read(item).decode().replace("</Relationships>", rel + "</Relationships>").encode()
            elif item == "[Content_Types].xml":
                content = zin.read(item).decode().replace("</Types>", override + "</Types>").encode()
            zout.writestr(item, content)
        zout.writestr("ppt/diagrams/data1.xml", data_xml)
    tmp.unlink()


if __name__ == "__main__":
    for name, fn in list(globals().items()):
        if name.startswith("build_") and callable(fn):
            print(f"Building {name}...")
            fn()
